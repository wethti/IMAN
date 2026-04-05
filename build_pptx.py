"""
build_pptx.py — Generate Peace Corridor PowerPoint from slide content.
Uses python-pptx to produce a 16:9 deck styled with IMAN brand colors.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy, os

# ── Brand colors ──────────────────────────────────────────────────────────────
BG       = RGBColor(0x0C, 0x1E, 0x0C)   # near-black green
BG2      = RGBColor(0x15, 0x2B, 0x15)   # dark green panel
GREEN    = RGBColor(0x7F, 0xBA, 0x00)   # lime accent
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE= RGBColor(0xDD, 0xEE, 0xDD)
MUTED    = RGBColor(0x7A, 0x8E, 0x7A)
NEWS_BG  = RGBColor(0x0A, 0x1F, 0x0A)

# ── Slide dimensions: 16:9 widescreen ─────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout

# ── Helper: add solid background ──────────────────────────────────────────────
def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

# ── Helper: add a filled rectangle ────────────────────────────────────────────
def rect(slide, l, t, w, h, fill_color=None, border_color=None, border_pt=1):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # 1 = MSO_SHAPE_TYPE.RECTANGLE
    if fill_color:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_pt)
    else:
        shape.line.fill.background()
    return shape

# ── Helper: add a text box ─────────────────────────────────────────────────────
def txbox(slide, l, t, w, h, text, size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, italic=False, font="Barlow", wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font
    return tb

# ── Helper: multi-paragraph text box ──────────────────────────────────────────
def multi_txbox(slide, l, t, w, h, paragraphs, wrap=True):
    """
    paragraphs = list of dicts:
      { text, size, bold, color, italic, font, align, space_before, space_after }
    """
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    first = True
    for p_def in paragraphs:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.alignment = p_def.get("align", PP_ALIGN.LEFT)
        if p_def.get("space_before"):
            p.space_before = Pt(p_def["space_before"])
        if p_def.get("space_after"):
            p.space_after  = Pt(p_def["space_after"])
        run = p.add_run()
        run.text           = p_def.get("text", "")
        run.font.size      = Pt(p_def.get("size", 16))
        run.font.bold      = p_def.get("bold", False)
        run.font.italic    = p_def.get("italic", False)
        run.font.color.rgb = p_def.get("color", WHITE)
        run.font.name      = p_def.get("font", "Barlow")
    return tb

# ── Helper: green accent bar (left edge) ──────────────────────────────────────
def green_bar(slide, top=Inches(1.2), height=Inches(5.1)):
    rect(slide, Inches(0.45), top, Pt(4), height, fill_color=GREEN)

# ── Helper: label (small caps-style) ──────────────────────────────────────────
def label(slide, l, t, text, size=9):
    txbox(slide, l, t, Inches(11), Inches(0.3), text,
          size=size, bold=True, color=GREEN, font="Barlow")

# ── Helper: photo placeholder ──────────────────────────────────────────────────
def photo_placeholder(slide, l, t, w, h, caption=""):
    rect(slide, l, t, w, h, fill_color=BG2, border_color=GREEN, border_pt=0.75)
    txbox(slide, l + Inches(0.2), t + h//2 - Inches(0.4),
          w - Inches(0.4), Inches(0.8),
          f"📷  {caption}", size=9, italic=True, color=MUTED,
          align=PP_ALIGN.CENTER, font="Barlow")

# ── Helper: bullet list ────────────────────────────────────────────────────────
def bullet_list(slide, l, t, w, h, items):
    """items = list of (bold_prefix, rest_of_text) tuples"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, (bold_part, rest) in enumerate(items):
        if first: p = tf.paragraphs[0]; first = False
        else:     p = tf.add_paragraph()
        p.space_before = Pt(4)
        # Bullet dot
        dot = p.add_run(); dot.text = "● "
        dot.font.color.rgb = GREEN; dot.font.size = Pt(9); dot.font.name = "Barlow"
        # Bold prefix
        if bold_part:
            b = p.add_run(); b.text = bold_part
            b.font.bold = True; b.font.size = Pt(13)
            b.font.color.rgb = WHITE; b.font.name = "Barlow"
        # Rest of text
        if rest:
            r = p.add_run(); r.text = rest
            r.font.size = Pt(13); r.font.color.rgb = OFF_WHITE; r.font.name = "Barlow"

# ── Helper: callout box ────────────────────────────────────────────────────────
def callout(slide, l, t, w, h, title, body):
    rect(slide, l, t, w, h, fill_color=RGBColor(0x15, 0x30, 0x15),
         border_color=GREEN, border_pt=0.75)
    txbox(slide, l + Inches(0.15), t + Inches(0.12),
          w - Inches(0.3), Inches(0.25),
          title, size=8, bold=True, color=GREEN, font="Barlow")
    txbox(slide, l + Inches(0.15), t + Inches(0.38),
          w - Inches(0.3), h - Inches(0.5),
          body, size=11, italic=True, color=OFF_WHITE, font="Barlow")

# ── Helper: section divider line ──────────────────────────────────────────────
def divider(slide, l, t, w=Inches(0.6)):
    rect(slide, l, t, w, Pt(3), fill_color=GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Right panel
rect(s, Inches(7.5), 0, Inches(5.833), H, fill_color=BG2)
photo_placeholder(s, Inches(7.5), 0, Inches(5.833), H,
                  "IMAN community engagement team, Englewood")

# IMAN badge
rect(s, Inches(0.5), Inches(0.55), Inches(4.2), Inches(0.32),
     fill_color=None, border_color=GREEN, border_pt=0.75)
txbox(s, Inches(0.55), Inches(0.58), Inches(4.1), Inches(0.28),
      "IMAN · INNER-CITY MUSLIM ACTION NETWORK",
      size=8, bold=True, color=GREEN, font="Barlow")

# Program name
txbox(s, Inches(0.45), Inches(1.0), Inches(6.8), Inches(3.5),
      "PEACE\nCORRIDOR™",
      size=88, bold=True, color=WHITE, font="Oswald")

# Tagline
txbox(s, Inches(0.5), Inches(4.7), Inches(6.5), Inches(0.7),
      "Street Outreach · Transitional Employment · Hospital Intervention\n"
      "A holistic approach to ending gun violence in Englewood.",
      size=13, color=OFF_WHITE, font="Barlow")

# Geography
txbox(s, Inches(0.5), Inches(5.6), Inches(3), Inches(0.35),
      "— ENGLEWOOD, CHICAGO",
      size=9, bold=True, color=MUTED, font="Barlow")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MISSION STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
divider(s, Inches(6.17), Inches(1.1))
multi_txbox(s, Inches(1.0), Inches(1.3), Inches(11.333), Inches(4.5), [
    {"text": "Peace Corridor ", "size": 26, "bold": True, "color": WHITE, "font": "Barlow"},
], wrap=True)

# Use a single large centered text box
tb = s.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(5.0))
tb.word_wrap = True
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER

runs = [
    ("Peace Corridor", True, WHITE, False),
    (" provides street-level violence interruption, transitional employment, "
     "cognitive behavioral therapy, and hospital-based crisis intervention for ", False, OFF_WHITE, False),
    ("highest-risk individuals and justice-involved youth", True, WHITE, False),
    (", as a holistic approach to ", False, OFF_WHITE, False),
    ("ending gun violence, preventing premature death, and restoring safety "
     "to Englewood's streets", True, GREEN, False),
    (".", False, OFF_WHITE, False),
]
for text, bold, color, italic in runs:
    r = p.add_run(); r.text = text
    r.font.size = Pt(22); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Barlow"

divider(s, Inches(6.17), Inches(6.4))
txbox(s, Inches(0), Inches(6.8), W, Inches(0.35),
      "IMAN  ·  PEACE CORRIDOR™",
      size=9, bold=True, color=MUTED, align=PP_ALIGN.CENTER, font="Barlow")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CORE PROGRAM TRACKS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
label(s, Inches(0.5), Inches(0.3), "THREE CORE PROGRAM TRACKS")

cards = [
    ("📢", "STREET OUTREACH\n& INTERRUPTION",
     "Credible messengers embedded in Englewood's highest-shooting blocks detect conflicts "
     "before they escalate, mediate disputes, and connect highest-risk individuals to "
     "services. Built on the Cure Violence model — proven to reduce shootings by 30–60%.",
     "FULL-TIME · COMMUNITY-EMBEDDED"),
    ("🧠", "CBT + TRANSITIONAL\nEMPLOYMENT",
     "18-month program combining cognitive behavioral therapy with paid transitional jobs, "
     "modeled after READI Chicago — shown to reduce violent crime involvement by 79% among "
     "highest-risk men. Direct pipeline into IMAN's Green ReEntry trades training.",
     "18 MONTHS · WEEKDAYS"),
    ("🏥", "HOSPITAL-BASED\nINTERVENTION",
     "Violence intervention specialists embedded in the trauma bay reach shooting victims "
     "at the moment of highest motivation for change — before discharge, before retaliation. "
     "A 30-minute conversation at that moment saves lives.",
     "24/7 · ON-CALL"),
]
card_w = Inches(4.0); gap = Inches(0.27)
for i, (icon, title, desc, chip) in enumerate(cards):
    l = Inches(0.5) + i * (card_w + gap)
    rect(s, l, Inches(0.7), card_w, Inches(5.4),
         fill_color=BG2, border_color=GREEN, border_pt=1)
    # Top accent bar
    rect(s, l, Inches(0.7), card_w, Pt(4), fill_color=GREEN)
    txbox(s, l + Inches(0.15), Inches(0.85), card_w - Inches(0.3), Inches(0.4),
          icon, size=22, font="Barlow")
    txbox(s, l + Inches(0.15), Inches(1.35), card_w - Inches(0.3), Inches(0.9),
          title, size=17, bold=True, color=WHITE, font="Oswald")
    txbox(s, l + Inches(0.15), Inches(2.35), card_w - Inches(0.3), Inches(2.9),
          desc, size=11.5, color=OFF_WHITE, font="Barlow")
    # Chip
    rect(s, l + Inches(0.15), Inches(5.55), card_w - Inches(0.3), Inches(0.38),
         fill_color=None, border_color=GREEN, border_pt=0.75)
    txbox(s, l + Inches(0.2), Inches(0.1) + Inches(5.55), card_w - Inches(0.5), Inches(0.28),
          chip, size=8, bold=True, color=GREEN, font="Barlow")

txbox(s, Inches(0), Inches(7.1), W, Inches(0.3),
      "IMAN  ·  PEACE CORRIDOR™", size=8, bold=True, color=MUTED,
      align=PP_ALIGN.RIGHT, font="Barlow")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WRAPAROUND SUPPORT SERVICES
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
label(s, Inches(0.5), Inches(0.3), "WRAPAROUND SUPPORT SERVICES")

supports = [
    ("LEGAL SERVICES",
     "Expungement, record sealing, and civil legal aid — removing the criminal "
     "record barriers that block housing, employment, and stability. Partnerships "
     "with clinics across family, housing, and civil rights law."),
    ("WORKFORCE PLACEMENT",
     "Direct pipeline into IMAN's Green ReEntry trades program. Graduates earn "
     "NCCER-accredited certifications in carpentry, welding, HVAC, and EV technology "
     "with living-wage job placement assistance."),
    ("ALUMNI LEADERSHIP",
     "Program graduates become paid Peace Corridor outreach workers and community "
     "safety ambassadors — turning lived experience into professional credibility "
     "and cycles into careers."),
]
card_w = Inches(4.0)
for i, (title, desc) in enumerate(supports):
    l = Inches(0.5) + i * (card_w + gap)
    rect(s, l, Inches(0.75), card_w, Inches(3.8),
         fill_color=None, border_color=GREEN, border_pt=0.75)
    rect(s, l, Inches(0.75), card_w, Pt(3), fill_color=GREEN)
    txbox(s, l + Inches(0.15), Inches(0.92), card_w - Inches(0.3), Inches(0.5),
          title, size=14, bold=True, color=GREEN, font="Oswald")
    txbox(s, l + Inches(0.15), Inches(1.52), card_w - Inches(0.3), Inches(2.8),
          desc, size=12, color=OFF_WHITE, font="Barlow")

# Callout bottom
callout(s, Inches(0.5), Inches(5.1), Inches(12.333), Inches(1.5),
        "NO HANDOFFS. ONE TEAM.",
        "Every participant moves through a single integrated pathway — from first contact "
        "on the block or in a trauma bay, to stable employment, housing, and legal standing.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — THEORY OF CHANGE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Left panel text
txbox(s, Inches(0.5), Inches(0.45), Inches(7.2), Inches(2.8),
      '"Everyone in Englewood wants to live.\nNot everyone has been given\na reason to believe they will."',
      size=28, bold=True, color=WHITE, font="Oswald")

divider(s, Inches(0.5), Inches(3.4))

txbox(s, Inches(0.5), Inches(3.65), Inches(7.0), Inches(2.4),
      "If individuals at the highest risk of being shot or incarcerated are reached by "
      "credible messengers, provided with economic alternatives through transitional "
      "employment and CBT, and connected to healthcare, housing, and legal support — "
      "then they will be on a thriving pathway to financial, physical, social-emotional, "
      "and spiritual growth that will impact themselves, their families, and their community.",
      size=13, color=OFF_WHITE, font="Barlow")

txbox(s, Inches(0.5), Inches(6.15), Inches(7.0), Inches(0.65),
      "The evidence is not theoretical. It is documented.\n"
      "What Englewood has been missing is not proof — it is investment.",
      size=12, bold=True, italic=True, color=GREEN, font="Barlow")

# Right panel
rect(s, Inches(8.0), 0, Inches(5.333), H, fill_color=BG2)
photo_placeholder(s, Inches(8.0), 0, Inches(5.333), H,
                  "IMAN community safety team member — in action on the block")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOLISTIC APPROACH
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Left vertical label bar
rect(s, 0, 0, Inches(0.85), H, fill_color=BG2)
txbox(s, -Inches(2.8), Inches(2.3), Inches(6.5), Inches(0.6),
      "A HOLISTIC APPROACH",
      size=14, bold=True, color=GREEN, font="Oswald")

sections = [
    ("Crisis-to-stability continuum for every participant",
     "From the moment of contact — whether on the block or in a hospital trauma bay — "
     "every Peace Corridor participant is enrolled in a structured pathway: outreach → CBT → "
     "employment → housing → healthcare → legal. No handoffs. One team."),
    ("Weekly therapy, financial literacy & behavioral health",
     "Every participant receives 1 hour of individual therapy and 1 hour of life coaching "
     "per week, plus weekly group sessions in financial literacy, trauma recovery, and "
     "spiritual grounding — gaining more than a job, gaining a foundation."),
    ("Full integration with IMAN's Health Center, Green ReEntry & Griot Plaza",
     "Peace Corridor participants are enrolled in IMAN's FQHC, connected to Green ReEntry's "
     "workforce pipeline, housed in transitional units, and anchored at Griot Plaza. "
     "This is the full IMAN ecosystem — activated for Englewood's most vulnerable."),
]
for i, (title, body) in enumerate(sections):
    t = Inches(0.7) + i * Inches(2.15)
    rect(s, Inches(0.95), t, Pt(3), Inches(1.8), fill_color=GREEN)
    txbox(s, Inches(1.2), t, Inches(7.5), Inches(0.4),
          title, size=12, bold=True, color=GREEN, font="Barlow")
    txbox(s, Inches(1.2), t + Inches(0.42), Inches(7.5), Inches(1.5),
          body, size=11.5, color=OFF_WHITE, font="Barlow")

# Ecosystem diagram (right side)
eco_l = Inches(9.4); eco_t = Inches(1.2)
eco_size = Inches(3.5)
rect(s, eco_l, eco_t, eco_size, eco_size, fill_color=None, border_color=GREEN, border_pt=0.75)
# Inner ring
rect(s, eco_l + Inches(0.8), eco_t + Inches(0.8),
     eco_size - Inches(1.6), eco_size - Inches(1.6),
     fill_color=None, border_color=RGBColor(0x50, 0x80, 0x20), border_pt=0.5)
txbox(s, eco_l + Inches(1.1), eco_t + Inches(1.35), Inches(1.3), Inches(0.8),
      "PEACE\nCORRIDOR", size=9, bold=True, color=WHITE,
      align=PP_ALIGN.CENTER, font="Oswald")

nodes = [
    ("Outreach",  eco_l + Inches(1.3),  eco_t - Inches(0.22)),
    ("CBT",       eco_l + eco_size,      eco_t + Inches(0.6)),
    ("Jobs",      eco_l + eco_size,      eco_t + Inches(2.2)),
    ("Housing",   eco_l + Inches(1.3),  eco_t + eco_size - Inches(0.12)),
    ("Health",    eco_l - Inches(0.55), eco_t + Inches(2.2)),
    ("Legal",     eco_l - Inches(0.55), eco_t + Inches(0.6)),
]
for text, nl, nt in nodes:
    txbox(s, nl, nt, Inches(0.9), Inches(0.28),
          text, size=8, bold=True, color=GREEN, font="Barlow")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EVIDENCE BASE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

label(s, Inches(0.5), Inches(0.3), "BACKED BY THE STRONGEST EVIDENCE IN PUBLIC HEALTH.")
txbox(s, Inches(0.5), Inches(0.65), Inches(7.5), Inches(1.1),
      "The Most Rigorously Tested Violence Intervention Models in the U.S.\n— Applied in Englewood.",
      size=26, bold=True, color=WHITE, font="Oswald")

evidence = [
    ("Street Outreach — Cure Violence",
     "30–60% reduction in shootings · 56% homicide drop in Baltimore",
     "Ransford et al., systematic review; Johns Hopkins evaluation"),
    ("CBT + Employment — READI Chicago",
     "79% reduction in violent crime involvement",
     "University of Chicago Crime Lab, 2023"),
    ("Hospital-Based Intervention",
     "3–5× more likely to engage when reached in trauma bay",
     "National Network of Hospital-Based Violence Intervention Programs"),
]
for i, (model, stat, cite) in enumerate(evidence):
    t = Inches(1.9) + i * Inches(1.65)
    rect(s, Inches(0.5), t, Pt(4), Inches(1.3), fill_color=GREEN)
    txbox(s, Inches(0.65), t, Inches(7.2), Inches(0.3),
          model, size=9, bold=True, color=GREEN, font="Barlow")
    txbox(s, Inches(0.65), t + Inches(0.3), Inches(7.2), Inches(0.55),
          stat, size=17, bold=True, color=WHITE, font="Barlow Condensed")
    txbox(s, Inches(0.65), t + Inches(0.88), Inches(7.2), Inches(0.3),
          cite, size=9, italic=True, color=MUTED, font="Barlow")

photo_placeholder(s, Inches(8.5), Inches(0.5), Inches(4.5), Inches(6.5),
                  "Participant in skills training or counseling session")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — STREET OUTREACH DEEP-DIVE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
txbox(s, Inches(0.5), Inches(0.35), Inches(8), Inches(1.5),
      "STREET OUTREACH\n& INTERRUPTION",
      size=48, bold=True, color=WHITE, font="Oswald")
bullet_list(s, Inches(0.5), Inches(2.05), Inches(12.333), Inches(3.2), [
    ("Cure Violence model proven in 30+ US cities", " — average 30–60% reduction in shootings "
     "when implemented with fidelity; 56% homicide reduction in Baltimore's highest-violence district"),
    ("Outreach workers are credible messengers", " — formerly justice-involved individuals from "
     "the same blocks, trained in conflict mediation, CBT techniques, and case management; "
     "the only messengers high-risk individuals trust"),
    ("Real-time intelligence layer", " — outreach workers in Englewood's highest-shooting "
     "corridors detect brewing conflicts before police or hospitals see them, coordinating "
     "with IMAN's existing community safety team"),
])
callout(s, Inches(0.5), Inches(5.6), Inches(12.333), Inches(1.3),
        "YEAR 1 PROJECTION",
        "6–10 full-time outreach workers covering Englewood's highest-risk corridors · "
        "Estimated 30–50% reduction in shootings in target area within 24 months")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — CBT + EMPLOYMENT DEEP-DIVE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
txbox(s, Inches(0.5), Inches(0.35), Inches(8), Inches(1.5),
      "CBT + TRANSITIONAL\nEMPLOYMENT",
      size=48, bold=True, color=WHITE, font="Oswald")
bullet_list(s, Inches(0.5), Inches(2.05), Inches(12.333), Inches(3.2), [
    ("79% reduction in violent crime involvement", " among READI Chicago participants — "
     "highest-risk men who received 18 months of transitional employment + CBT vs. control group "
     "(UChicago Crime Lab, 2023)"),
    ("Direct pipeline into Green ReEntry", " — participants earn NCCER-accredited certifications "
     "in carpentry, welding, HVAC, and EV technology with living-wage job placement assistance"),
    ("91% retention rate in 2024", " across 350 Green ReEntry participants — proving IMAN's "
     "workforce model holds even for the highest-risk populations"),
])
callout(s, Inches(0.5), Inches(5.6), Inches(12.333), Inches(1.3),
        "SUCCESS STORY",
        "One Green ReEntry graduate is now a certified instructor at City Colleges of Chicago "
        "— teaching the next cohort the same skills that changed his life.")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — LEGAL SERVICES DEEP-DIVE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

txbox(s, Inches(0.5), Inches(0.35), Inches(7.5), Inches(1.5),
      "LEGAL\nSERVICES",
      size=48, bold=True, color=WHITE, font="Oswald")
bullet_list(s, Inches(0.5), Inches(2.05), Inches(7.5), Inches(3.0), [
    ("Criminal records are the wall", " between graduates and stability — expungement and "
     "record sealing remove the single biggest structural barrier to housing, employment, "
     "and financial access"),
    ("Partnerships with legal clinics and firms", " across family law, housing law, driving "
     "records, and civil rights — the same infrastructure IMAN uses in Green ReEntry, now "
     "extended to Peace Corridor"),
    ("Know Your Rights workshops", " — every participant understands their rights in police "
     "encounters, housing disputes, and employment — because knowledge is the first form "
     "of protection"),
])
txbox(s, Inches(0.5), Inches(5.5), Inches(7.5), Inches(0.75),
      "Green ReEntry has secured expungements for 20+ participants and hosted Know Your Rights "
      "workshops for nearly 50. Peace Corridor extends this to 50+ participants annually in Year 1.",
      size=10, italic=True, color=MUTED, font="Barlow")

photo_placeholder(s, Inches(8.3), Inches(0.5), Inches(4.7), Inches(6.5),
                  "IMAN participants at the Illinois State Capitol")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FOCUS ON HEALING
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)
txbox(s, Inches(0.5), Inches(0.35), Inches(8), Inches(1.5),
      "A FOCUS ON\nHEALING",
      size=48, bold=True, color=WHITE, font="Oswald")
txbox(s, Inches(0.5), Inches(1.95), Inches(12.333), Inches(0.55),
      "As part of IMAN's holistic model, all Peace Corridor participants receive regular "
      "primary and behavioral health care through IMAN's Federally Qualified Health Center:",
      size=13, bold=True, color=OFF_WHITE, font="Barlow")
bullet_list(s, Inches(0.5), Inches(2.65), Inches(12.333), Inches(2.5), [
    ("IMAN's Behavioral Health team has seen over 2,000 patients from Green ReEntry",
     " — the same integrated care infrastructure, already proven, now extended to "
     "Peace Corridor's highest-risk population"),
    ("Specialized trauma recovery groups", " — trauma recovery, grief counseling, parenting "
     "skills, and spiritual care — offered in community, not in a clinic"),
])
txbox(s, Inches(0.5), Inches(5.55), Inches(12.333), Inches(0.75),
      "Gun violence is not only a public safety crisis. It is a public health crisis.\n"
      "Peace Corridor treats it as both.",
      size=14, bold=True, italic=True, color=GREEN, font="Barlow")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — PROJECTED HIGHLIGHTS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Vertical green accent bar
rect(s, Inches(0.45), Inches(0.55), Pt(5), Inches(6.4), fill_color=GREEN)

txbox(s, Inches(0.8), Inches(0.45), Inches(10), Inches(1.1),
      "YEAR 1–2 PROJECTED HIGHLIGHTS",
      size=42, bold=True, color=WHITE, font="Oswald")

stats = [
    ("50–75",    "highest-risk individuals",
     "enrolled in CBT + transitional employment in Year 1, drawn from Englewood's "
     "documented highest-risk population (projected)"),
    ("200+",     "conflict mediations",
     "conducted by street outreach team in Year 1 — each one a potential life saved "
     "before it becomes a shooting statistic (projected)"),
    ("30–50%",   "projected reduction in shootings",
     "in target corridors within 24 months — consistent with Cure Violence outcomes "
     "in comparable cities (projected)"),
]
for i, (num, bold_desc, rest_desc) in enumerate(stats):
    t = Inches(1.65) + i * Inches(1.75)
    txbox(s, Inches(0.8), t, Inches(3.5), Inches(0.95),
          num, size=52, bold=True, color=GREEN, font="Oswald")
    tb = s.shapes.add_textbox(Inches(4.2), t + Inches(0.1), Inches(8.7), Inches(0.9))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = bold_desc + " "
    r1.font.size = Pt(14); r1.font.bold = True
    r1.font.color.rgb = WHITE; r1.font.name = "Barlow"
    r2 = p.add_run(); r2.text = rest_desc
    r2.font.size = Pt(13); r2.font.color.rgb = OFF_WHITE; r2.font.name = "Barlow"

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — BUILDING THE FUTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

# Vertical label
rect(s, 0, 0, Inches(0.85), H, fill_color=BG2)
txbox(s, -Inches(3.2), Inches(2.0), Inches(7.5), Inches(0.7),
      "BUILDING THE FUTURE",
      size=14, bold=True, color=GREEN, font="Oswald")

cols = [
    ("PROGRAM SCALE", [
        ("Grow to 100-Person Cohorts",
         "Double the CBT + employment cohort as the Regenerator opens — providing "
         "on-site housing, pharmacy, and FQHC for the program's highest-need participants"),
        ("Expand to Chicago Lawn",
         "Extend street outreach and hospital-based intervention to IMAN's second primary "
         "service neighborhood, building a South Side-wide safety infrastructure"),
    ]),
    ("WORKFORCE INTEGRATION", [
        ("Culinary, Automotive & EV Trades",
         "As Green ReEntry expands into three new trades tracks, Peace Corridor graduates "
         "gain additional pathways including the growing EV infrastructure sector"),
        ("Formalize the Alumni Pipeline",
         "Structured career track: participant → graduate → paid outreach worker → "
         "program staff — turning lived experience into a profession"),
    ]),
    ("HOSPITAL PARTNERSHIPS", [
        ("2–3 South Side Trauma Centers",
         "Expand hospital-based intervention to reach shooting victims at every major "
         "South Side trauma bay — including U of C Medical Center and Advocate Christ"),
    ]),
]
col_w = Inches(3.8)
for i, (cat, items) in enumerate(cols):
    l = Inches(1.05) + i * (col_w + Inches(0.22))
    # Category header
    rect(s, l, Inches(0.45), col_w, Inches(0.35),
         fill_color=None, border_color=GREEN, border_pt=0.75)
    txbox(s, l + Inches(0.1), Inches(0.5), col_w - Inches(0.2), Inches(0.28),
          cat, size=9, bold=True, color=WHITE, font="Barlow")
    for j, (title, body) in enumerate(items):
        t = Inches(1.05) + j * Inches(2.95)
        txbox(s, l, t, col_w, Inches(0.38),
              title, size=12, bold=True, color=GREEN, font="Barlow")
        txbox(s, l, t + Inches(0.4), col_w, Inches(2.3),
              body, size=11, color=OFF_WHITE, font="Barlow")
    # Column divider
    if i < 2:
        rect(s, l + col_w + Inches(0.08), Inches(0.45),
             Pt(1), Inches(6.6), fill_color=RGBColor(0x2A, 0x4A, 0x2A))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — IN THE NEWS
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s, NEWS_BG)

# Header box
rect(s, Inches(0.5), Inches(0.35), Inches(4.5), Inches(0.65),
     fill_color=None, border_color=GREEN, border_pt=1.5)
txbox(s, Inches(0.6), Inches(0.4), Inches(4.3), Inches(0.55),
      "IN THE NEWS:", size=26, bold=True, color=WHITE, font="Oswald")

news = [
    ("CHICAGO TRIBUNE",
     "'I always thought I'd be a statistic': South Side program helps men rebuild after prison"),
    ("CBS CHICAGO",
     "'Green ReEntry' Program Seeks To Reduce Recidivism By Getting Formerly Incarcerated Work Experience"),
    ("NPR",
     "Green ReEntry Program Takes Holistic Approach To Reducing Recidivism"),
    ("NEXT CITY",
     "A Safe House in Chicago Provides More Than Shelter"),
]
for i, (outlet, headline) in enumerate(news):
    t = Inches(1.2) + i * Inches(1.4)
    txbox(s, Inches(0.5), t, Inches(12.333), Inches(0.32),
          outlet, size=10, bold=True, color=GREEN, font="Barlow")
    txbox(s, Inches(0.5), t + Inches(0.32), Inches(12.333), Inches(0.7),
          headline, size=13, color=OFF_WHITE, font="Barlow")
    rect(s, Inches(0.5), t + Inches(1.05), Inches(12.333), Pt(0.5),
         fill_color=RGBColor(0x2A, 0x4A, 0x2A))

txbox(s, Inches(0.5), Inches(6.9), Inches(12.333), Inches(0.45),
      "Peace Corridor builds on IMAN's nationally recognized track record — bringing the same "
      "credibility, the same infrastructure, and the same trust to Englewood's most urgent unmet need.",
      size=11, italic=True, color=MUTED, font="Barlow")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — ABOUT IMAN
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s)

txbox(s, Inches(0.5), Inches(0.35), Inches(7.5), Inches(1.1),
      "ABOUT IMAN:", size=52, bold=True, color=WHITE, font="Oswald")
txbox(s, Inches(0.5), Inches(1.45), Inches(7.5), Inches(0.32),
      "FOSTERING HEALTH, WELLNESS, AND HEALING",
      size=10, bold=True, color=GREEN, font="Barlow")
divider(s, Inches(0.5), Inches(1.85))

paras = [
    ("The ", False, OFF_WHITE), ("Inner-City Muslim Action Network", True, WHITE),
    (" (IMAN) is a community organization that fosters health, wellness, and healing "
     "in the inner city by organizing for social change, cultivating the arts, and "
     "operating a holistic health center.", False, OFF_WHITE),
]
tb = s.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(7.4), Inches(1.1))
tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
for text, bold, color in paras:
    r = p.add_run(); r.text = text
    r.font.size = Pt(12.5); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Barlow"

txbox(s, Inches(0.5), Inches(3.35), Inches(7.4), Inches(1.1),
      "IMAN incorporated as a nonprofit in 1997 through the drive of people directly "
      "affected by and deeply invested in social issues affecting communities of color "
      "living on Chicago's South Side. Since that time, IMAN has steadily grown and, "
      "in 2016, opened a second office in Atlanta.",
      size=12.5, color=OFF_WHITE, font="Barlow")

txbox(s, Inches(0.5), Inches(4.6), Inches(7.4), Inches(1.4),
      "The organization models an integrative approach incorporating primary and "
      "behavioral health; artistic expression; leadership development; organizing and "
      "advocacy; housing; and job training — in an effort to substantially increase "
      "the quality of life for people in marginalized communities.",
      size=12.5, color=OFF_WHITE, font="Barlow")

photo_placeholder(s, Inches(8.3), 0, Inches(5.033), H,
                  "IMAN community march and organizing action")

# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════
out = "/workspaces/IMAN/Data library/Peace Corridor Pitch Deck.pptx"
prs.save(out)
print(f"✓ Saved: {out}")
print(f"  Slides: {len(prs.slides)}")
