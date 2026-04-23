"""Build C2_pitch_deck.pptx using python-pptx, mirroring the visual language
of IMAN Research Progress v4.pptx (10.00 x 5.625, Calibri, dark green header,
lime accent, F4F7F1 page background)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(os.path.dirname(HERE), "deliverable", "C2_pitch_deck.pptx")

# ===== Palette (from Research Progress v4) =====
BG         = RGBColor(0xF4, 0xF7, 0xF1)  # page
INK        = RGBColor(0x1A, 0x1A, 0x1A)
GREEN_DARK = RGBColor(0x1A, 0x3A, 0x1A)  # header bar, headline
GREEN_2    = RGBColor(0x16, 0x2E, 0x16)  # deep card
GREEN_MID  = RGBColor(0x2D, 0x5A, 0x1B)
LIME       = RGBColor(0x7F, 0xBA, 0x00)  # accent strip, key accent
CARD_LIGHT = RGBColor(0xEB, 0xF2, 0xE8)  # card fill
CARD_2     = RGBColor(0xF0, 0xF5, 0xED)
SAGE_MUTE  = RGBColor(0xAA, 0xBB, 0xA0)
MUTED      = RGBColor(0x5A, 0x6B, 0x5A)  # eyebrow + small caps
TITLE_TXT  = RGBColor(0xFF, 0xFF, 0xFF)
PAGE_NUM   = RGBColor(0x4A, 0x7A, 0x4A)
FOOTER_TXT = RGBColor(0x8C, 0xAA, 0x8C)
SUBTL_PALE = RGBColor(0xC8, 0xDC, 0xC8)
TEXT_BODY  = RGBColor(0x26, 0x2E, 0x22)
TERRACOTTA = RGBColor(0xB8, 0x50, 0x42)
AMBER      = RGBColor(0xE4, 0xB2, 0x5A)
GRID_LN    = RGBColor(0xC8, 0xD2, 0xC6)

FOOTER_STR = "IMAN 2040  •  Investor Brief  •  April 2026"

# ===== Helpers =====

def new_deck():
    p = Presentation()
    p.slide_width = Inches(10.0)
    p.slide_height = Inches(5.625)
    return p


def blank_slide(p):
    layout = p.slide_layouts[6]  # Blank
    s = p.slides.add_slide(layout)
    # Page background
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, p.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.shadow.inherit = False
    return s


def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def rect(slide, x, y, w, h, rgb, no_line=True):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    _fill(sh, rgb)
    if no_line:
        _no_line(sh)
    sh.shadow.inherit = False
    return sh


def hline(slide, x, y, w, rgb, weight=0.75):
    sh = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    sh.line.color.rgb = rgb
    sh.line.width = Pt(weight)
    return sh


def text_box(slide, x, y, w, h, text, *, font="Calibri", size=10, bold=False,
             italic=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=None, space_after=0, margin=0.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if isinstance(text, str):
        parts = [text]
    else:
        parts = list(text)
    for i, seg in enumerate(parts):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(seg, dict):
            body = seg.get("text", "")
            r = p.add_run()
            r.text = body
            r.font.name = seg.get("font", font)
            r.font.size = Pt(seg.get("size", size))
            r.font.bold = seg.get("bold", bold)
            r.font.italic = seg.get("italic", italic)
            r.font.color.rgb = seg.get("color", color)
        else:
            r = p.add_run()
            r.text = seg
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
    return tb


def multi_para(slide, x, y, w, h, paragraphs, *, align=PP_ALIGN.LEFT,
               anchor=MSO_ANCHOR.TOP, margin=0.05, line_spacing=None):
    """Each item in paragraphs is a dict with text, font, size, bold, color, bullet, space_after, indent."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    for i, seg in enumerate(paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = seg.get("align", align)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if seg.get("space_after"):
            p.space_after = Pt(seg["space_after"])
        runs = seg.get("runs")
        if runs:
            for rdef in runs:
                r = p.add_run()
                r.text = rdef.get("text", "")
                r.font.name = rdef.get("font", "Calibri")
                r.font.size = Pt(rdef.get("size", 10))
                r.font.bold = rdef.get("bold", False)
                r.font.italic = rdef.get("italic", False)
                r.font.color.rgb = rdef.get("color", INK)
        else:
            r = p.add_run()
            body = seg.get("text", "")
            if seg.get("bullet"):
                body = "•  " + body
            r.text = body
            r.font.name = seg.get("font", "Calibri")
            r.font.size = Pt(seg.get("size", 10))
            r.font.bold = seg.get("bold", False)
            r.font.italic = seg.get("italic", False)
            r.font.color.rgb = seg.get("color", INK)
        if seg.get("indent"):
            p.level = seg["indent"]
    return tb


# ===== Chrome =====

def header(slide, title, page_num):
    rect(slide, 0, 0, 10.0, 0.85, GREEN_DARK)
    rect(slide, 0, 0.85, 10.0, 0.06, LIME)
    text_box(slide, 0.45, 0.09, 8.30, 0.71, title,
             font="Calibri", size=21, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    text_box(slide, 9.00, 0.09, 0.90, 0.71, page_num,
             font="Calibri", size=10, color=SUBTL_PALE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)


def footer(slide):
    text_box(slide, 0.45, 5.34, 9.10, 0.22, FOOTER_STR,
             font="Calibri", size=8.5, color=FOOTER_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)


def eyebrow(slide, text, y=1.08):
    text_box(slide, 0.45, y, 9.10, 0.24, text.upper(),
             font="Calibri", size=9, bold=True, color=MUTED,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)


def content_slide(p, title, page_num, eyebrow_text=None):
    s = blank_slide(p)
    header(s, title, page_num)
    if eyebrow_text:
        eyebrow(s, eyebrow_text)
    footer(s)
    return s


# ===== SLIDE BUILDERS =====

def slide_01_cover(p):
    """Cover — dark, quiet, numbers-first."""
    s = p.slides.add_slide(p.slide_layouts[6])
    # Full dark bg
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, p.slide_height)
    _fill(bg, GREEN_DARK); _no_line(bg); bg.shadow.inherit = False
    # Lime vertical strip
    rect(s, 0, 0, 0.28, 5.625, LIME)
    # Eyebrow
    text_box(s, 0.65, 0.55, 8.5, 0.30, "INNER-CITY MUSLIM ACTION NETWORK  —  ENGLEWOOD 2040",
             font="Calibri", size=10, color=LIME, margin=0.0)
    # Title
    text_box(s, 0.65, 1.00, 8.8, 0.85, "The Englewood Question",
             font="Calibri", size=40, bold=True, color=TITLE_TXT, margin=0.0)
    # Subtitle
    multi_para(s, 0.65, 1.95, 8.8, 1.10, [
        {"runs": [{"text": "Why a 30-year life expectancy gap exists, why it hasn’t closed,",
                   "size": 18, "color": SUBTL_PALE}]},
        {"runs": [{"text": "and what would actually move it.",
                   "size": 18, "color": SUBTL_PALE}]},
    ], margin=0.0)
    # Accent rule
    rect(s, 0.65, 3.30, 3.5, 0.04, LIME)
    # Location / date
    text_box(s, 0.65, 3.48, 8.8, 0.32, "Chicago, IL   •   April 2026   •   Evidence-based briefing",
             font="Calibri", size=12, color=SUBTL_PALE, margin=0.0)
    # Hero stat pair (ticker bottom)
    # Numbers: bottom-anchored so the baseline sits above the caption.
    text_box(s, 0.65, 4.15, 3.00, 0.80, "90",
             font="Calibri", size=34, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.BOTTOM, margin=0.0)
    text_box(s, 0.65, 5.02, 3.00, 0.28, "Life expectancy, Streeterville",
             font="Calibri", size=10, color=SUBTL_PALE, margin=0.0)
    text_box(s, 3.85, 4.15, 3.00, 0.80, "67.7",
             font="Calibri", size=34, bold=True, color=LIME,
             anchor=MSO_ANCHOR.BOTTOM, margin=0.0)
    text_box(s, 3.85, 5.02, 3.00, 0.28, "Life expectancy, Englewood",
             font="Calibri", size=10, color=SUBTL_PALE, margin=0.0)
    text_box(s, 7.05, 4.15, 2.50, 0.80, "− 22.3",
             font="Calibri", size=34, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.BOTTOM, margin=0.0)
    text_box(s, 7.05, 5.02, 2.50, 0.28, "Years of life, same city",
             font="Calibri", size=10, color=SUBTL_PALE, margin=0.0)
    return s


def slide_02_one_sentence(p):
    s = content_slide(p, "22 years of life, in the same city",
                      "02 / 18", eyebrow_text="The one sentence")
    # Large pull quote (left) + comparison block (right)
    multi_para(s, 0.45, 1.45, 5.40, 3.40, [
        {"runs": [{"text": "The gap is not new. It has been measured ",
                   "size": 18, "color": TEXT_BODY},
                  {"text": "by a different name in every generation",
                   "size": 18, "color": GREEN_DARK, "bold": True}]},
        {"space_after": 10, "runs": [{"text": " ", "size": 8}]},
        {"runs": [{"text": "“health disparities” in the 1980s, “structural racism” in the 2000s, ",
                   "size": 12, "color": MUTED},
                  {"text": "“social determinants of health” today. The label rotates.",
                   "size": 12, "color": MUTED}]},
        {"space_after": 8, "runs": [{"text": " ", "size": 8}]},
        {"runs": [{"text": "The number has moved very little.",
                   "size": 13, "color": GREEN_DARK, "bold": True, "italic": True}]},
    ], line_spacing=1.15)
    # Comparison card (right)
    rect(s, 6.10, 1.40, 3.45, 3.40, CARD_LIGHT)
    text_box(s, 6.30, 1.55, 3.15, 0.28,
             "LIFE EXPECTANCY BY CHICAGO TRACT",
             font="Calibri", size=9, bold=True, color=MUTED, margin=0.0)
    rows = [
        ("Streeterville",  "90.0", GREEN_DARK),
        ("Lincoln Park",   "85.1", GREEN_MID),
        ("Hyde Park",      "80.4", GREEN_MID),
        ("Chicago average", "77.3", MUTED),
        ("Little Village", "82.0", GREEN_MID),
        ("Chicago Lawn",   "73.9", MUTED),
        ("Englewood",      "67.7", TERRACOTTA),
    ]
    for i, (name, val, col) in enumerate(rows):
        y = 1.95 + i * 0.38
        text_box(s, 6.30, y, 2.00, 0.32, name,
                 font="Calibri", size=10.5, color=TEXT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        text_box(s, 8.30, y, 1.15, 0.32, val,
                 font="Calibri", size=13, bold=True, color=col,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    return s


def slide_03_complication(p):
    s = content_slide(p, "Two decades of programs. A neighborhood that hasn’t moved.",
                      "03 / 18", eyebrow_text="The complication")
    # Intro paragraph
    multi_para(s, 0.45, 1.45, 9.10, 0.95, [
        {"runs": [{"text": "Englewood has been a destination for programs. Grocery stores have opened. Clinics have expanded. Violence-interruption crews have patrolled. Some of these programs have been excellent — and rigorously evaluated. ",
                   "size": 11.5, "color": TEXT_BODY},
                  {"text": "And yet Englewood’s life expectancy has barely budged.",
                   "size": 11.5, "color": GREEN_DARK, "bold": True}]},
    ], line_spacing=1.2)
    # 3-column "proven program" row
    card_y = 2.55
    card_h = 1.55
    cards = [
        ("Barbershop BP program",
         "−27 mmHg",
         "systolic blood pressure reduction",
         "LA / Cedars-Sinai RCT"),
        ("Cure Violence / READI",
         "−50%",
         "serious violence over 18 months",
         "Chicago West Side"),
        ("Health-center expansion",
         "+1.5 yr",
         "life expectancy in service area",
         "FQHC panel evidence"),
    ]
    col_w = 2.95
    gap = 0.12
    start_x = 0.45
    for i, (name, val, caption, src) in enumerate(cards):
        x = start_x + i * (col_w + gap)
        rect(s, x, card_y, col_w, card_h, CARD_LIGHT)
        rect(s, x, card_y, col_w, 0.04, LIME)
        text_box(s, x + 0.15, card_y + 0.12, col_w - 0.3, 0.34, name.upper(),
                 font="Calibri", size=9, bold=True, color=MUTED, margin=0.0)
        text_box(s, x + 0.15, card_y + 0.42, col_w - 0.3, 0.55, val,
                 font="Calibri", size=26, bold=True, color=GREEN_DARK, margin=0.0)
        text_box(s, x + 0.15, card_y + 0.98, col_w - 0.3, 0.30, caption,
                 font="Calibri", size=10, color=TEXT_BODY, margin=0.0)
        text_box(s, x + 0.15, card_y + 1.28, col_w - 0.3, 0.22, src,
                 font="Calibri", size=8.5, italic=True, color=MUTED, margin=0.0)
    # Impact line
    text_box(s, 0.45, 4.40, 9.10, 0.55,
             "The programs work. The neighborhood does not change. Life expectancy is a neighborhood outcome, not a program outcome.",
             font="Calibri", size=11.5, bold=True, italic=True, color=GREEN_DARK,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    return s


def slide_04_variance(p):
    s = content_slide(p, "The remaining variance is about place, not behavior",
                      "04 / 18", eyebrow_text="What the data actually says")
    # Left: explanation text
    multi_para(s, 0.45, 1.40, 3.55, 3.60, [
        {"runs": [{"text": "Hold behavior constant. The gap is still there.",
                   "size": 11, "color": TEXT_BODY}]},
        {"space_after": 6, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "Hold income constant. The gap is still there.",
                   "size": 11, "color": TEXT_BODY}]},
        {"space_after": 6, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "Hold education and insurance constant. The gap shrinks but doesn’t close.",
                   "size": 11, "color": TEXT_BODY}]},
        {"space_after": 14, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "Roughly half of tract-level life-expectancy variance in the U.S. is structural / place-based ",
                   "size": 11, "color": GREEN_DARK, "bold": True},
                  {"text": "— not individual behavior, not healthcare access, not income alone. That is the share IMAN is set up to move.",
                   "size": 11, "color": GREEN_DARK, "bold": True}]},
    ], line_spacing=1.25)
    # Right: chart image
    chart_path = os.path.join(HERE, "chart_variance.png")
    s.shapes.add_picture(chart_path, Inches(4.25), Inches(1.40), width=Inches(5.30))
    # Caption under chart
    text_box(s, 4.25, 5.02, 5.30, 0.22,
             "Decomposition of tract-level life-expectancy variance, pooled U.S. data (CDC USALEEP; peer-reviewed models, 2018–2024).",
             font="Calibri", size=8.5, italic=True, color=MUTED, margin=0.0)
    return s


def slide_05_deviants(p):
    s = content_slide(p, "Some poor neighborhoods outlive their demographics by a decade",
                      "05 / 18", eyebrow_text="Positive deviants")
    # Chart on left
    chart_path = os.path.join(HERE, "chart_scatter.png")
    s.shapes.add_picture(chart_path, Inches(0.40), Inches(1.30), width=Inches(5.05))
    # Right column: what's happening
    multi_para(s, 5.55, 1.45, 4.10, 3.60, [
        {"runs": [{"text": "Each dot is a U.S. census tract.",
                   "size": 10, "italic": True, "color": MUTED}]},
        {"space_after": 10, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "The dashed line is what the model predicts given poverty rate.",
                   "size": 10.5, "color": TEXT_BODY}]},
        {"space_after": 10, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "The amber dots — Little Village, East Boston, Sunset Park, Chinatown — sit 10–15 years above their predicted line. ",
                   "size": 10.5, "color": TEXT_BODY},
                  {"text": "They are poor. They are dense. They are non-white. Their people live longer than the model predicts.",
                   "size": 10.5, "color": GREEN_DARK, "bold": True}]},
        {"space_after": 10, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "Englewood (red) sits roughly on the line. ",
                   "size": 10.5, "color": TEXT_BODY},
                  {"text": "The deviants prove the gap is not fate.",
                   "size": 10.5, "color": TERRACOTTA, "bold": True}]},
    ], line_spacing=1.2)
    return s


def slide_06_collective_efficacy(p):
    s = content_slide(p, "The feature in common is not a program. It’s a form.",
                      "06 / 18", eyebrow_text="Collective efficacy — Sampson, 25 years of evidence")
    # Four feature cards
    feats = [
        ("Walkable streets",
         "Daily pedestrian encounters outrank gym membership as a cardiovascular predictor."),
        ("Three-generation households",
         "Social support that hospitals price at $30K / yr per chronic-disease patient. Free, if it exists."),
        ("Anchor institutions",
         "A mosque, church, market, school, or cultural center that residents encounter daily."),
        ("Dense networks",
         "People who know each other’s children’s names — the single most consistent non-behavioral predictor of who lives and dies."),
    ]
    col_w = 2.22
    gap = 0.12
    x0 = 0.45
    y0 = 1.45
    card_h = 2.55
    for i, (name, body) in enumerate(feats):
        x = x0 + i * (col_w + gap)
        rect(s, x, y0, col_w, card_h, CARD_LIGHT)
        rect(s, x, y0, col_w, 0.60, GREEN_DARK)
        text_box(s, x + 0.15, y0 + 0.12, col_w - 0.3, 0.36, name,
                 font="Calibri", size=11.5, bold=True, color=TITLE_TXT,
                 anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        text_box(s, x + 0.15, y0 + 0.75, col_w - 0.3, card_h - 0.9, body,
                 font="Calibri", size=10, color=TEXT_BODY,
                 line_spacing=1.2, margin=0.0)
    # Bottom takeaway bar
    rect(s, 0.45, 4.30, 9.10, 0.62, GREEN_DARK)
    text_box(s, 0.65, 4.30, 8.80, 0.62,
             "IMAN’s goal is not to import clinical interventions into Englewood. It is to manufacture the structural features of a naturally long-lived poor neighborhood around the people who already live there.",
             font="Calibri", size=11, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    return s


def slide_07_stack_thesis(p):
    s = content_slide(p, "Single programs move inches. Stacked ecosystems move years.",
                      "07 / 18", eyebrow_text="The stacked-intervention thesis")
    # Chart (left) + explanation (right)
    chart_path = os.path.join(HERE, "chart_stacked.png")
    s.shapes.add_picture(chart_path, Inches(0.40), Inches(1.30), width=Inches(5.35))
    multi_para(s, 5.85, 1.45, 3.80, 3.55, [
        {"runs": [{"text": "At full dose in the target population:",
                   "size": 10, "italic": True, "color": MUTED}]},
        {"space_after": 10, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "~7 years gross LE gain,",
                   "size": 18, "bold": True, "color": GREEN_DARK}]},
        {"runs": [{"text": " summed across IMAN’s seven pillars.",
                   "size": 11, "color": TEXT_BODY}]},
        {"space_after": 10, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "Net — after overlap, attrition, selection — the defensible projection is ",
                   "size": 10.5, "color": TEXT_BODY},
                  {"text": "3–5 years over 10–15 years",
                   "size": 10.5, "bold": True, "color": GREEN_DARK},
                  {"text": " in a majority-Black neighborhood at full build-out.",
                   "size": 10.5, "color": TEXT_BODY}]},
        {"space_after": 8, "runs": [{"text": " ", "size": 4}]},
        {"runs": [{"text": "Dudley Street, East Lake, Purpose Built — wherever the stack has been run, effects compounded non-linearly. Adjacent programs reinforce each other’s retention and dose.",
                   "size": 10, "color": MUTED, "italic": True}]},
    ], line_spacing=1.18)
    return s


def slide_08_chicago_gap(p):
    s = content_slide(p, "Same city. Same hospitals. Same mayor. 22 years apart.",
                      "08 / 18", eyebrow_text="Chicago, at tract resolution")
    # Left: pair of vertical columns (Streeterville vs Englewood) as tall bars
    # Frame
    rect(s, 0.45, 1.40, 3.55, 3.55, CARD_LIGHT)
    # Small caption at top of card (above bars)
    text_box(s, 0.55, 1.52, 3.40, 0.22,
             "YEARS OF LIFE EXPECTANCY",
             font="Calibri", size=8.5, bold=True, color=MUTED,
             align=PP_ALIGN.CENTER, margin=0.0)
    # Bars — bar_base sits well above the card bottom so the city labels fit.
    bar_base = 4.35
    bar_top_max = 2.05
    # Streeterville
    rect(s, 1.10, bar_top_max, 0.80, bar_base - bar_top_max, GREEN_DARK)
    text_box(s, 0.85, bar_top_max - 0.34, 1.30, 0.32, "90.0",
             font="Calibri", size=16, bold=True, color=GREEN_DARK,
             align=PP_ALIGN.CENTER, margin=0.0)
    text_box(s, 0.85, bar_base + 0.04, 1.30, 0.30, "Streeterville",
             font="Calibri", size=9.5, color=TEXT_BODY,
             align=PP_ALIGN.CENTER, margin=0.0)
    # Englewood
    eng_h = (67.7 / 90.0) * (bar_base - bar_top_max)
    eng_top = bar_base - eng_h
    rect(s, 2.50, eng_top, 0.80, eng_h, TERRACOTTA)
    text_box(s, 2.25, eng_top - 0.34, 1.30, 0.32, "67.7",
             font="Calibri", size=16, bold=True, color=TERRACOTTA,
             align=PP_ALIGN.CENTER, margin=0.0)
    text_box(s, 2.25, bar_base + 0.04, 1.30, 0.30, "Englewood",
             font="Calibri", size=9.5, color=TEXT_BODY,
             align=PP_ALIGN.CENTER, margin=0.0)
    # Gap note sitting below the city-name row, still inside the card.
    text_box(s, 0.55, 4.70, 3.40, 0.22,
             "a 22.3-year gap, 3 miles apart",
             font="Calibri", size=9, italic=True, color=GREEN_DARK,
             align=PP_ALIGN.CENTER, margin=0.0)
    # Right: list of "if you normalize by X, gap becomes Y" rows
    text_box(s, 4.30, 1.40, 5.30, 0.32, "THE GAP IS ROBUST TO ADJUSTMENT",
             font="Calibri", size=9, bold=True, color=MUTED, margin=0.0)
    adj = [
        ("Raw crude gap",                     "22.3 yr"),
        ("− poverty rate",               "14.1 yr"),
        ("− income + education",         "10.8 yr"),
        ("− insurance coverage",         "9.4 yr"),
        ("− race / ethnicity",           "8.1 yr"),
        ("− all of the above",           "≈ 5 yr unexplained"),
    ]
    row_y = 1.78
    for label, val in adj:
        color_val = GREEN_DARK if "unexplained" in val else TEXT_BODY
        text_box(s, 4.30, row_y, 3.20, 0.32, label,
                 font="Calibri", size=10.5, color=TEXT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        text_box(s, 7.55, row_y, 2.05, 0.32, val,
                 font="Calibri", size=11.5, bold=True, color=color_val,
                 align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        hline(s, 4.30, row_y + 0.44, 5.30, GRID_LN, weight=0.5)
        row_y += 0.48
    text_box(s, 4.30, 4.70, 5.30, 0.28,
             "Even after every statistical adjustment, about 5 years of gap remain. That residual is structural.",
             font="Calibri", size=9.5, italic=True, color=GREEN_DARK, margin=0.0)
    return s


def slide_09_ecosystem(p):
    s = content_slide(p, "IMAN already operates a full-ecosystem footprint in Englewood",
                      "09 / 18", eyebrow_text="What is already on the ground")
    # Six pillars as a 3x2 grid
    pillars = [
        ("IMAN Health Center",  "10,000+ patients; primary care, behavioral health, addiction medicine."),
        ("Fresh Beat / Fresh Market", "Halal grocery + corner-store conversions on 63rd Street corridor."),
        ("Green ReEntry",       "Housing rehab + construction workforce for returning citizens."),
        ("Arts & Culture",      "Griot Plaza, open-air venues, cultural programming as civic anchor."),
        ("Community Organizing","Weekly civic formation; leadership pipelines from Chicago Lawn to Englewood."),
        ("Regenerator / Policy","Civic-tech platform connecting residents to benefits, jobs, services."),
    ]
    col_w = 2.95
    gap = 0.12
    row_h = 1.55
    y0 = 1.40
    x0 = 0.45
    for i, (name, desc) in enumerate(pillars):
        col = i % 3
        row = i // 3
        x = x0 + col * (col_w + gap)
        y = y0 + row * (row_h + 0.12)
        rect(s, x, y, col_w, row_h, CARD_LIGHT)
        rect(s, x, y, 0.10, row_h, LIME)
        text_box(s, x + 0.25, y + 0.12, col_w - 0.35, 0.40, name,
                 font="Calibri", size=12, bold=True, color=GREEN_DARK,
                 anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        text_box(s, x + 0.25, y + 0.55, col_w - 0.35, row_h - 0.70, desc,
                 font="Calibri", size=10, color=TEXT_BODY,
                 line_spacing=1.2, margin=0.0)
    # Takeaway
    text_box(s, 0.45, 4.80, 9.10, 0.28,
             "The 2040 ask is not to start — it is to intensify, to the density at which stacked effects appear.",
             font="Calibri", size=10.5, italic=True, bold=True, color=GREEN_DARK,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    return s


def pillar_slide(p, num, title, page, eye, headline_val, headline_caption,
                 interventions, evidence_src, impact_line):
    """Shared layout for pillar slides 10-15."""
    s = content_slide(p, title, page, eyebrow_text=eye)
    # Left big-number panel
    rect(s, 0.45, 1.40, 3.10, 3.50, GREEN_DARK)
    rect(s, 0.45, 1.40, 3.10, 0.06, LIME)
    text_box(s, 0.60, 1.60, 2.80, 0.30, f"PILLAR {num} — AT FULL DOSE",
             font="Calibri", size=9, bold=True, color=LIME, margin=0.0)
    text_box(s, 0.60, 1.95, 2.80, 1.10, headline_val,
             font="Calibri", size=54, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    text_box(s, 0.60, 3.05, 2.80, 0.60, headline_caption,
             font="Calibri", size=10.5, color=SUBTL_PALE,
             line_spacing=1.2, margin=0.0)
    text_box(s, 0.60, 4.55, 2.80, 0.30, evidence_src,
             font="Calibri", size=8.5, italic=True, color=SAGE_MUTE, margin=0.0)
    # Right side: what IMAN does
    rect(s, 3.75, 1.40, 5.80, 3.50, CARD_LIGHT)
    text_box(s, 3.95, 1.55, 5.40, 0.30, "WHAT IMAN ACTUALLY DOES",
             font="Calibri", size=9, bold=True, color=MUTED, margin=0.0)
    row_y = 1.95
    for bullet in interventions:
        # bullet marker
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.00), Inches(row_y + 0.10),
                                 Inches(0.10), Inches(0.10))
        _fill(dot, LIME); _no_line(dot); dot.shadow.inherit = False
        text_box(s, 4.20, row_y, 5.15, 0.45, bullet,
                 font="Calibri", size=10.5, color=TEXT_BODY,
                 line_spacing=1.2, margin=0.0)
        row_y += 0.45
    # Impact strip at bottom of right card
    rect(s, 3.75, 4.50, 5.80, 0.40, GREEN_DARK)
    text_box(s, 3.95, 4.50, 5.40, 0.40, impact_line,
             font="Calibri", size=10.5, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    return s


def slide_10_health_center(p):
    return pillar_slide(p, 1,
        "Health Center expansion",
        "10 / 18",
        "Pillar 1 — clinical base",
        "+1.5", "years of LE gain in the FQHC service area at full panel penetration.",
        [
            "Expand to 20,000 primary-care patients on the South Side.",
            "Integrate behavioral health, addiction medicine, pharmacy on-site.",
            "Embedded chronic-disease coaching: BP, A1c, lipid management.",
            "Link every intake to food-access, housing, and workforce referrals.",
        ],
        "FQHC evidence panel (HRSA / NEJM, 2016–2024).",
        "Full-panel FQHC evidence sits at +1.5 yr of LE gain — the anchor of the ecosystem.")


def slide_11_green_reentry(p):
    return pillar_slide(p, 2,
        "Green ReEntry — housing plus workforce",
        "11 / 18",
        "Pillar 2 — structural",
        "+2.0", "years of LE gain in the housed-and-employed cohort (Dudley Street analog).",
        [
            "Rehab 250 vacant homes in Englewood / Chicago Lawn over 10 years.",
            "Construction workforce staffed from returning-citizen cohort.",
            "Living-wage jobs that build the neighborhood they live in.",
            "Housing stability eliminates the chief driver of missed-care cycles.",
        ],
        "Dudley Street Neighborhood Initiative; returning-citizen housing RCTs.",
        "Housing + income stability is the single largest structural lever in the stack.")


def slide_12_fresh_market(p):
    return pillar_slide(p, 3,
        "Fresh Market + clinical coaching",
        "12 / 18",
        "Pillar 3 — food access done right",
        "+0.5", "years of LE gain — contingent on clinical wraparound.",
        [
            "Halal grocery + corner-store conversions across 63rd Street.",
            "Produce prescriptions paired with Health Center coaching.",
            "Fresh Beat SNAP-match: subsidized purchases tied to appointments.",
            "NOTE: produce prescriptions alone have ~zero effect — IMAN’s model pairs them with clinical dose.",
        ],
        "Philadelphia Healthy Corner Store; 2026 produce-Rx pooled trials.",
        "Food access only moves outcomes when bundled with the clinical pillar. IMAN’s model is bundled.")


def slide_13_violence(p):
    return pillar_slide(p, 4,
        "Violence interruption",
        "13 / 18",
        "Pillar 4 — community safety",
        "+0.7", "years of LE gain in young-adult men at READI-level dose.",
        [
            "Formal partnership with Cure Violence / READI Chicago.",
            "Hire 8–12 interrupters from the Green ReEntry pool.",
            "Weekly shooting reviews with Health Center behavioral-health team.",
            "Three-year minimum funding commitment — shorter cycles erase the effect.",
        ],
        "Cure Violence RCTs; READI Chicago evaluation (UChicago, 2022–2024).",
        "READI-dose violence interruption cuts serious violence 40–65% — the clearest causal effect in the stack.")


def slide_14_barbershop(p):
    return pillar_slide(p, 5,
        "Barbershop blood-pressure program",
        "14 / 18",
        "Pillar 5 — cardiovascular beachhead",
        "+0.8", "years of LE gain in the Black-male population at full reach.",
        [
            "Deploy in 8–10 Black-owned barbershops on 63rd & Halsted.",
            "IMAN pharmacist prescribes in-shop; barbers screen during cuts.",
            "$25 voucher per completed visit drives repeat adherence.",
            "64% reach BP control vs. 12% in usual care (Cedars-Sinai RCT).",
        ],
        "Victor et al., NEJM 2018; replication studies 2021–2024.",
        "One of the largest treatment effects ever measured in cardiology. Deployable in year 1.")


def slide_15_griot(p):
    return pillar_slide(p, 6,
        "Griot Plaza and community organizing",
        "15 / 18",
        "Pillar 6 — the collective-efficacy pillar",
        "×", "this pillar is the multiplier that makes the other five compound, rather than stack arithmetically.",
        [
            "Open-air plaza + performance venue as daily neighborhood anchor.",
            "Weekly civic formation: leadership pipelines, youth programming.",
            "Arts programming that brings 20,000+ residents through the block/yr.",
            "Mosque + Health Center + Plaza + Market together = the anchoring ecology Sampson documents.",
        ],
        "Sampson, Great American City; 25 years of Chicago PHDCN data.",
        "There is no “program” for collective efficacy — only the cumulative effect of dense, durable institutions.")


def slide_16_honest_projection(p):
    s = content_slide(p, "What to promise — and what not to",
                      "16 / 18", eyebrow_text="The defensible projection")
    # Large central headline
    rect(s, 0.45, 1.40, 9.10, 1.40, GREEN_DARK)
    text_box(s, 0.70, 1.55, 8.60, 0.40, "AT FULL BUILD-OUT, OVER 10–15 YEARS",
             font="Calibri", size=10, bold=True, color=LIME, margin=0.0)
    text_box(s, 0.70, 1.90, 8.60, 0.85,
             "3 to 5 years of life-expectancy gain in a majority-Black neighborhood.",
             font="Calibri", size=22, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    # Two columns: promise / don't promise
    rect(s, 0.45, 3.00, 4.45, 1.95, CARD_LIGHT)
    rect(s, 0.45, 3.00, 4.45, 0.36, GREEN_MID)
    text_box(s, 0.60, 3.00, 4.15, 0.36, "WHY NOT MORE",
             font="Calibri", size=9.5, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    multi_para(s, 0.60, 3.45, 4.15, 1.45, [
        {"runs": [{"text": "7–10 year numbers come from immigrant enclaves whose longevity reflects culturally-bounded diet, immigration selection, and the ‘returning-home’ effect that statistically inflates measured longevity. Structural features transfer. Demographic tailwinds do not.",
                   "size": 10, "color": TEXT_BODY}]}
    ], line_spacing=1.2)
    rect(s, 5.10, 3.00, 4.45, 1.95, CARD_LIGHT)
    rect(s, 5.10, 3.00, 4.45, 0.36, GREEN_MID)
    text_box(s, 5.25, 3.00, 4.15, 0.36, "WHY NOT LESS",
             font="Calibri", size=9.5, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    multi_para(s, 5.25, 3.45, 4.15, 1.45, [
        {"runs": [{"text": "Stacked-intervention literature, adjacent evidence from Little Village, and dose-response curves from every single-intervention RCT point in the same direction. 3–5 years is the floor of the defensible range, not a ceiling.",
                   "size": 10, "color": TEXT_BODY}]}
    ], line_spacing=1.2)
    return s


def slide_17_measurement(p):
    s = content_slide(p, "Three nested layers. Three audiences. One number that matters.",
                      "17 / 18", eyebrow_text="How we will know")
    layers = [
        ("QUARTERLY DASHBOARD", "For the board",
         "BP control, ED visits for ambulatory-care-sensitive conditions, fresh-produce share of Fresh Market sales, housing placements, monthly homicides & shootings in the service area.",
         "moves in months"),
        ("SYNTHETIC CONTROL", "For funders",
         "From 2028: is Englewood diverging from its 16 peer communities (South Side Chicago, West Detroit, Memphis, Jackson, East St. Louis, Flint, Youngstown, …)? Direction + magnitude of divergence separates IMAN from macro drift.",
         "moves in years"),
        ("LE RESIDUAL", "For the historical record",
         "When the successor to USALEEP ships (≈2030): did Englewood’s actual life expectancy end up higher than its poverty, race, insurance, and unemployment profile predicted? That is the number the work will be remembered by.",
         "moves in a decade"),
    ]
    x = 0.45
    y = 1.40
    col_w = 2.95
    gap = 0.12
    card_h = 3.60
    for i, (head, aud, body, tempo) in enumerate(layers):
        cx = x + i * (col_w + gap)
        rect(s, cx, y, col_w, card_h, CARD_LIGHT)
        rect(s, cx, y, col_w, 0.55, GREEN_DARK)
        text_box(s, cx + 0.15, y + 0.10, col_w - 0.3, 0.35, head,
                 font="Calibri", size=10, bold=True, color=LIME,
                 anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
        text_box(s, cx + 0.15, y + 0.70, col_w - 0.3, 0.30, aud.upper(),
                 font="Calibri", size=9, bold=True, color=MUTED, margin=0.0)
        text_box(s, cx + 0.15, y + 1.05, col_w - 0.3, 2.15, body,
                 font="Calibri", size=10, color=TEXT_BODY, line_spacing=1.25, margin=0.0)
        # Tempo badge
        rect(s, cx + 0.15, y + card_h - 0.55, col_w - 0.30, 0.40, GREEN_DARK)
        text_box(s, cx + 0.15, y + card_h - 0.55, col_w - 0.30, 0.40, tempo,
                 font="Calibri", size=10, bold=True, color=LIME,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    text_box(s, 0.45, 5.08, 9.10, 0.22,
             "Leading indicators every quarter. Lagging outcome once a decade. The board doesn’t wait for the historical record.",
             font="Calibri", size=9.5, italic=True, color=MUTED, margin=0.0)
    return s


def slide_18_ask(p):
    s = p.slides.add_slide(p.slide_layouts[6])
    # Dark stage
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, p.slide_width, p.slide_height)
    _fill(bg, GREEN_DARK); _no_line(bg); bg.shadow.inherit = False
    # Left vertical lime strip
    rect(s, 0, 0, 0.28, 5.625, LIME)
    # Eyebrow
    text_box(s, 0.65, 0.55, 8.5, 0.30, "THE ASK",
             font="Calibri", size=10, bold=True, color=LIME, margin=0.0)
    # Big number
    text_box(s, 0.65, 1.00, 6.0, 1.60, "$25M",
             font="Calibri", size=96, bold=True, color=TITLE_TXT,
             anchor=MSO_ANCHOR.MIDDLE, margin=0.0)
    text_box(s, 0.65, 2.60, 8.5, 0.45, "catalytic capital to intensify the ecosystem to stacked-evidence density.",
             font="Calibri", size=16, color=SUBTL_PALE, margin=0.0)
    # Thin rule
    rect(s, 0.65, 3.10, 3.5, 0.04, LIME)
    # Four line items
    items = [
        "Fund the ecosystem, not the program.",
        "Give it ten to fifteen years, not three.",
        "Measure the lagging outcome, but watch the leading indicators every quarter.",
        "Compare Englewood to the communities it actually resembles.",
    ]
    y = 3.30
    for line in items:
        # dot
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), Inches(y + 0.09),
                                 Inches(0.10), Inches(0.10))
        _fill(dot, LIME); _no_line(dot); dot.shadow.inherit = False
        text_box(s, 0.90, y, 8.5, 0.32, line,
                 font="Calibri", size=13, color=TITLE_TXT, margin=0.0)
        y += 0.38
    # Close attribution
    text_box(s, 0.65, 5.20, 8.5, 0.30,
             "Inner-City Muslim Action Network   •   Englewood, Chicago   •   April 2026",
             font="Calibri", size=10, color=SUBTL_PALE, margin=0.0)
    return s


# ===== Driver =====

def main():
    p = new_deck()
    slide_01_cover(p)
    slide_02_one_sentence(p)
    slide_03_complication(p)
    slide_04_variance(p)
    slide_05_deviants(p)
    slide_06_collective_efficacy(p)
    slide_07_stack_thesis(p)
    slide_08_chicago_gap(p)
    slide_09_ecosystem(p)
    slide_10_health_center(p)
    slide_11_green_reentry(p)
    slide_12_fresh_market(p)
    slide_13_violence(p)
    slide_14_barbershop(p)
    slide_15_griot(p)
    slide_16_honest_projection(p)
    slide_17_measurement(p)
    slide_18_ask(p)
    p.save(OUT)
    print("wrote", OUT)


if __name__ == "__main__":
    main()
